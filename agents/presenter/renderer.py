from __future__ import annotations

import re
from pathlib import Path

import structlog
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from agents.presenter.outline import DeckOutline
from agents.presenter.slide_writer import SlideContent

log = structlog.get_logger(__name__)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Layout indices in a standard blank template
LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_BLANK = 6

# Scrim overlay: 40% opaque black over bottom third of image slides
SCRIM_ALPHA = "40000"
SCRIM_TOP = Inches(4.5)
SCRIM_HEIGHT = Inches(3.0)


def _slug(text: str) -> str:
    return re.sub(r"[^\w]+", "_", text.lower()).strip("_")[:40]


def _add_image_slide(
    prs: Presentation,
    content: SlideContent,
    image_path: Path | None,
) -> None:
    layout = prs.slide_layouts[LAYOUT_BLANK]
    slide = prs.slides.add_slide(layout)

    # Full-bleed image behind content
    if image_path and image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            left=Inches(0),
            top=Inches(0),
            width=SLIDE_W,
            height=SLIDE_H,
        )

    # Dark scrim at bottom third so text is readable over any background
    scrim = slide.shapes.add_shape(1, Inches(0), SCRIM_TOP, SLIDE_W, SCRIM_HEIGHT)
    scrim.fill.solid()
    scrim.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    scrim.line.fill.background()
    srgb = scrim.fill._fill._solidFill.find(qn("a:srgbClr"))
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", SCRIM_ALPHA)

    # Text box on top of scrim
    txBox = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(5.2),
        width=Inches(12.3),
        height=Inches(2.0),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Headline
    headline_para = tf.paragraphs[0]
    headline_para.text = content.headline
    run = headline_para.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Bullets
    for bullet in content.body:
        para = tf.add_paragraph()
        para.text = f"• {bullet}"
        run = para.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = content.speaker_note

    log.info(
        "renderer.slide.added", number=content.slide_number, type=content.slide_type
    )


def _add_text_slide(prs: Presentation, content: SlideContent) -> None:
    """For data/code_demo/team slides — text-only layout."""
    try:
        layout = prs.slide_layouts[LAYOUT_CONTENT]
    except IndexError:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # Try to use placeholder titles/content if available
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    if 0 in placeholders:
        placeholders[0].text = content.headline
    else:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)
        )
        run = txBox.text_frame.paragraphs[0]
        run.text = content.headline
        run.runs[0].font.size = Pt(32)
        run.runs[0].font.bold = True

    if 1 in placeholders:
        tf = placeholders[1].text_frame
        tf.text = ""
        for i, bullet in enumerate(content.body):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = bullet
            para.runs[0].font.size = Pt(18)
    else:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(content.body):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"• {bullet}"
            if para.runs:
                para.runs[0].font.size = Pt(18)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = content.speaker_note

    log.info(
        "renderer.slide.text", number=content.slide_number, type=content.slide_type
    )


def _add_title_slide(prs: Presentation, title: str) -> None:
    layout = prs.slide_layouts[LAYOUT_TITLE]
    slide = prs.slides.add_slide(layout)
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in placeholders:
        placeholders[0].text = title
    log.info("renderer.title_slide.added", title=title)


def render_deck(
    outline: DeckOutline,
    slides: list[SlideContent],
    image_map: dict[int, Path],
    template_path: Path | None,
    output_dir: Path,
) -> Path:
    """Build and save the .pptx. Returns the output path."""
    output_dir.mkdir(parents=True, exist_ok=True)

    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
        log.info("renderer.template.loaded", path=str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        log.info("renderer.template.blank")

    NO_IMAGE_TYPES = {"data", "code_demo", "team"}

    for content in slides:
        if content.slide_type in NO_IMAGE_TYPES:
            _add_text_slide(prs, content)
        else:
            image_path = image_map.get(content.slide_number)
            if image_path is None:
                log.warning(
                    "renderer.image.missing.fallback",
                    slide=content.slide_number,
                )
                _add_text_slide(prs, content)
            else:
                _add_image_slide(prs, content, image_path)

    deck_slug = _slug(outline.title)
    out_path = output_dir / f"{deck_slug}.pptx"
    prs.save(str(out_path))
    log.info("renderer.saved", path=str(out_path))
    return out_path
