"""Tests for slide renderer — scrim overlay and layout correctness."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches

from agents.presenter.models import DeckOutline, SlideContent, SlideOutline
from agents.presenter.renderer import SCRIM_ALPHA, render_deck


def _make_png(tmp_path: Path, name: str = "slide.png") -> Path:
    """Create a minimal valid PNG file for testing."""
    img = Image.new("RGB", (16, 16), color=(50, 100, 200))
    path = tmp_path / name
    img.save(path)
    return path


def _make_outline(title: str = "Test Deck") -> DeckOutline:
    return DeckOutline(
        title=title,
        slides=[
            SlideOutline(
                number=1,
                title="Intro",
                type="narrative",
                talking_points=["point a"],
                speaker_note="note",
            ),
            SlideOutline(
                number=2,
                title="Data",
                type="data",
                talking_points=["numbers"],
                speaker_note="note2",
            ),
        ],
    )


def _make_slides() -> list[SlideContent]:
    return [
        SlideContent(
            slide_number=1,
            slide_type="narrative",
            headline="The Intro",
            body=["bullet one", "bullet two"],
            speaker_note="presenter note",
            image_brief="An inspiring landscape",
        ),
        SlideContent(
            slide_number=2,
            slide_type="data",
            headline="The Numbers",
            body=["stat one"],
            speaker_note="data note",
            image_brief=None,
        ),
    ]


def test_render_deck_slide_count(tmp_path: Path) -> None:
    """render_deck produces one slide per SlideContent (no title slide added)."""
    outline = _make_outline()
    slides = _make_slides()
    out = render_deck(outline, slides, {}, None, tmp_path)

    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_image_slide_has_scrim(tmp_path: Path) -> None:
    """Image slide has ≥3 shapes: background image, scrim rectangle, text box."""
    img_path = _make_png(tmp_path)
    outline = _make_outline()
    slides = _make_slides()
    out = render_deck(outline, slides, {1: img_path}, None, tmp_path)

    prs = Presentation(str(out))
    image_slide = prs.slides[0]  # slide_number=1, type narrative (image slide)
    assert len(image_slide.shapes) >= 3  # image + scrim + textbox


def test_image_slide_scrim_is_black_with_alpha(tmp_path: Path) -> None:
    """The scrim rectangle on image slides has black fill with partial opacity."""
    img_path = _make_png(tmp_path)
    outline = _make_outline()
    slides = _make_slides()
    out = render_deck(outline, slides, {1: img_path}, None, tmp_path)

    prs = Presentation(str(out))
    image_slide = prs.slides[0]

    # Find a shape with solid black fill (the scrim)
    scrim_found = False
    for shape in image_slide.shapes:
        try:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is None:
                continue
            solid_fill = spPr.find(qn("a:solidFill"))
            if solid_fill is None:
                continue
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is None:
                continue
            if srgb.get("val", "").upper() == "000000":
                alpha_elem = srgb.find(qn("a:alpha"))
                if alpha_elem is not None:
                    assert alpha_elem.get("val") == SCRIM_ALPHA
                    scrim_found = True
                    break
        except (AttributeError, TypeError):
            continue

    assert scrim_found, "No scrim shape with black + alpha fill found on image slide"


def test_text_slide_has_no_scrim(tmp_path: Path) -> None:
    """Text-only slides (data/code_demo/team) are rendered without a scrim shape."""
    outline = _make_outline()
    slides = _make_slides()
    out = render_deck(outline, slides, {}, None, tmp_path)

    prs = Presentation(str(out))
    text_slide = prs.slides[1]  # slide_number=2, type data (text-only)

    # No shape should have a black solid fill with alpha (scrim)
    for shape in text_slide.shapes:
        try:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is None:
                continue
            solid_fill = spPr.find(qn("a:solidFill"))
            if solid_fill is None:
                continue
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None and srgb.get("val", "").upper() == "000000":
                alpha_elem = srgb.find(qn("a:alpha"))
                assert alpha_elem is None, "Text-only slide should not have a scrim"
        except (AttributeError, TypeError):
            continue


def test_image_slide_text_is_white(tmp_path: Path) -> None:
    """Headline text on image slides is white (readable over dark scrim)."""
    img_path = _make_png(tmp_path)
    outline = _make_outline()
    slides = _make_slides()
    out = render_deck(outline, slides, {1: img_path}, None, tmp_path)

    prs = Presentation(str(out))
    image_slide = prs.slides[0]

    # Find a text box with the headline
    headline_found = False
    for shape in image_slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "The Intro" in (run.text or ""):
                    assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
                    headline_found = True

    assert headline_found, "Headline text not found on image slide"
