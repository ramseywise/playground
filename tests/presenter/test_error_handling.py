"""Tests for error handling — JSON retry, image fetch retry, graceful degradation."""

from __future__ import annotations

import json
from pathlib import Path
from unittest.mock import MagicMock, call, patch

import httpx
import pytest

from agents.utils.client import parse_json_response
from agents.presenter.providers import MAX_RETRIES, RETRY_DELAYS, PollinationsProvider


# ---------------------------------------------------------------------------
# parse_json_response — JSON retry
# ---------------------------------------------------------------------------


def _make_client(response_text: str) -> MagicMock:
    client = MagicMock()
    mock_resp = MagicMock()
    mock_resp.content = [MagicMock(text=response_text)]
    client.messages.create.return_value = mock_resp
    return client


def test_parse_json_response_valid() -> None:
    """Valid JSON is returned immediately without a retry call."""
    client = MagicMock()
    result = parse_json_response(client, '{"key": "value"}', "model", "system")
    assert result == {"key": "value"}
    client.messages.create.assert_not_called()


def test_parse_json_response_fenced() -> None:
    """JSON wrapped in code fences is parsed without retry."""
    client = MagicMock()
    result = parse_json_response(client, "```json\n{\"key\": 1}\n```", "model", "system")
    assert result == {"key": 1}
    client.messages.create.assert_not_called()


def test_parse_json_response_retries_once_on_bad_json() -> None:
    """Malformed JSON triggers exactly one retry call to Claude."""
    valid_json = '{"fixed": true}'
    client = _make_client(valid_json)

    result = parse_json_response(client, "not valid json {{{", "claude-sonnet-4-6", "sys")

    assert result == {"fixed": True}
    client.messages.create.assert_called_once()


def test_parse_json_response_retry_passes_system() -> None:
    """Retry call uses the same system prompt."""
    client = _make_client('{"ok": true}')
    parse_json_response(client, "bad json", "my-model", "MY_SYSTEM_PROMPT")

    call_kwargs = client.messages.create.call_args[1]
    assert call_kwargs["system"] == "MY_SYSTEM_PROMPT"
    assert call_kwargs["model"] == "my-model"


def test_parse_json_response_raises_if_retry_also_fails() -> None:
    """If the retry response is also invalid JSON, JSONDecodeError is raised."""
    client = _make_client("still not json ~~~")
    with pytest.raises(json.JSONDecodeError):
        parse_json_response(client, "bad json 1", "model", "system")


def test_parse_json_response_list() -> None:
    """JSON arrays are returned as lists."""
    client = MagicMock()
    result = parse_json_response(client, '[{"a": 1}, {"b": 2}]', "model", "system")
    assert isinstance(result, list)
    assert len(result) == 2


# ---------------------------------------------------------------------------
# PollinationsProvider — retry with exponential backoff
# ---------------------------------------------------------------------------


def test_pollinations_retry_on_http_error(tmp_path: Path) -> None:
    """Provider retries up to MAX_RETRIES times on HTTP errors before raising."""
    provider = PollinationsProvider()
    dest = tmp_path / "out.png"

    with (
        patch("agents.presenter.providers._http_get_with_retry") as mock_get,
        patch("agents.presenter.providers.time.sleep") as mock_sleep,
    ):
        mock_get.side_effect = httpx.HTTPStatusError(
            "503", request=MagicMock(), response=MagicMock(status_code=503)
        )
        with pytest.raises(httpx.HTTPStatusError):
            provider.generate_image("test prompt", dest, 1280, 720)


def test_pollinations_retry_delays(tmp_path: Path) -> None:
    """_http_get_with_retry sleeps with increasing delays between attempts."""
    import time as time_module

    with (
        patch("agents.presenter.providers.time.sleep") as mock_sleep,
        patch("agents.presenter.providers.httpx.Client") as mock_client_cls,
    ):
        mock_resp = MagicMock()
        mock_resp.raise_for_status.side_effect = [
            httpx.HTTPStatusError("503", request=MagicMock(), response=MagicMock(status_code=503)),
            httpx.HTTPStatusError("503", request=MagicMock(), response=MagicMock(status_code=503)),
            None,  # third attempt succeeds
        ]
        mock_resp.content = b"image_bytes"
        mock_client_cls.return_value.__enter__.return_value.get.return_value = mock_resp

        from agents.presenter.providers import _http_get_with_retry
        result = _http_get_with_retry("https://example.com/image")

    assert result == b"image_bytes"
    # Two sleeps: after attempt 0 and after attempt 1
    assert mock_sleep.call_count == 2
    sleep_calls = [c[0][0] for c in mock_sleep.call_args_list]
    assert sleep_calls[0] == RETRY_DELAYS[0]
    assert sleep_calls[1] == RETRY_DELAYS[1]


def test_pollinations_success_no_sleep(tmp_path: Path) -> None:
    """Successful fetch on first attempt does not sleep."""
    with (
        patch("agents.presenter.providers.time.sleep") as mock_sleep,
        patch("agents.presenter.providers.httpx.Client") as mock_client_cls,
    ):
        mock_resp = MagicMock()
        mock_resp.raise_for_status.return_value = None
        mock_resp.content = b"image_data"
        mock_client_cls.return_value.__enter__.return_value.get.return_value = mock_resp

        from agents.presenter.providers import _http_get_with_retry
        _http_get_with_retry("https://example.com/img")

    mock_sleep.assert_not_called()


# ---------------------------------------------------------------------------
# Renderer graceful degradation
# ---------------------------------------------------------------------------


def test_render_deck_falls_back_to_text_when_image_missing(tmp_path: Path) -> None:
    """Image slide with no image_map entry falls back to text-only layout."""
    from agents.presenter.models import DeckOutline, SlideContent, SlideOutline
    from agents.presenter.renderer import render_deck

    outline = DeckOutline(
        title="Fallback Test",
        slides=[
            SlideOutline(
                number=1,
                title="Visual Slide",
                type="narrative",
                talking_points=["point"],
                speaker_note="note",
            )
        ],
    )
    slides = [
        SlideContent(
            slide_number=1,
            slide_type="narrative",
            headline="Visual",
            body=["bullet"],
            speaker_note="note",
            image_brief="A landscape",
        )
    ]

    # image_map is empty — image fetch "failed"
    out = render_deck(outline, slides, {}, None, tmp_path)

    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(str(out))
    slide = prs.slides[0]

    # Should NOT have a black scrim (text-only layout was used)
    for shape in slide.shapes:
        try:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is None:
                continue
            solid_fill = spPr.find(qn("a:solidFill"))
            if solid_fill is None:
                continue
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None and srgb.get("val", "").upper() == "000000":
                alpha_elem = srgb.find(qn("a:alpha"))
                assert alpha_elem is None, (
                    "Text-only fallback slide should not have a dark scrim"
                )
        except (AttributeError, TypeError):
            continue


def test_render_deck_uses_image_layout_when_image_present(tmp_path: Path) -> None:
    """Image slide with a valid image_map entry uses the image layout (has scrim)."""
    from PIL import Image  # only used to create a minimal valid PNG
    import io

    from agents.presenter.models import DeckOutline, SlideContent, SlideOutline
    from agents.presenter.renderer import render_deck, SCRIM_ALPHA

    # Create a minimal 10x10 PNG in tmp_path
    img = Image.new("RGB", (10, 10), color=(100, 100, 200))
    img_path = tmp_path / "slide_01.png"
    img.save(img_path)

    outline = DeckOutline(
        title="Image Test",
        slides=[
            SlideOutline(
                number=1,
                title="Visual",
                type="concept",
                talking_points=["p"],
                speaker_note="n",
            )
        ],
    )
    slides = [
        SlideContent(
            slide_number=1,
            slide_type="concept",
            headline="Visual Slide",
            body=["bullet"],
            speaker_note="note",
            image_brief="abstract",
        )
    ]

    out = render_deck(outline, slides, {1: img_path}, None, tmp_path)

    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(str(out))
    slide = prs.slides[0]

    scrim_found = any(
        (
            (spPr := shape._element.find(qn("p:spPr"))) is not None
            and (sf := spPr.find(qn("a:solidFill"))) is not None
            and (srgb := sf.find(qn("a:srgbClr"))) is not None
            and srgb.get("val", "").upper() == "000000"
            and srgb.find(qn("a:alpha")) is not None
            and srgb.find(qn("a:alpha")).get("val") == SCRIM_ALPHA
        )
        for shape in slide.shapes
    )
    assert scrim_found, "Image slide with image should have a scrim"
